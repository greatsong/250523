# file: survey_gpt_suite.py  (Streamlit 클라우드 메인 스크립트)

import streamlit as st
import pandas as pd
import numpy as np
import koreanize_matplotlib  # 사용자 필수
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
from io import StringIO
import json, re, asyncio, textwrap
from collections import Counter
from openai import OpenAI
from sklearn.cluster import KMeans
from sklearn.metrics import pairwise_distances
from sklearn.preprocessing import normalize
from sklearn.decomposition import PCA
from pptx import Presentation
from pptx.util import Inches, Pt
import pandasql as ps

# ───────────────────────── 설정 ───────────────────────── #
st.set_page_config("GPT 설문 슈트", "🧩", layout="wide")
client = OpenAI(api_key=st.secrets["openai_api_key"])
MODEL_CHAT = "gpt-4o-mini"
MODEL_EMB  = "text-embedding-3-large"

COLUMN_TYPES = {
    "timestamp":"타임스탬프","email":"이메일","phone":"전화번호","name":"이름",
    "numeric":"숫자","single_choice":"단일 선택","multiple_choice":"다중 선택",
    "text_short":"단답 텍스트","text_long":"장문 텍스트","other":"기타"
}
STOPWORDS = {'은','는','이','가','을','를','의','에','와','과','도','로','으로','만',
             '에서','까지','부터','라고','하고'}

# ────────────────────── GPT 유틸 함수 ────────────────────── #
@st.cache_data(show_spinner=False)
def gpt_infer_types(sample_csv:str)->dict:
    """GPT‑4o가 컬럼 의미 추론 → JSON"""
    res = client.chat.completions.create(
        model=MODEL_CHAT,
        messages=[
          {"role":"system","content":
           "You are a data scientist. Infer the semantic data type for each CSV column.\n"
           "Possible types: timestamp, email, phone, name, numeric, single_choice, multiple_choice, "
           "text_short, text_long, other. Return JSON."},
          {"role":"user","content":sample_csv}
        ],
        response_format={"type":"json_object"},
        temperature=0
    )
    return json.loads(res.choices[0].message.content)

def gpt_short_completion(sys,user,format_json=False):
    """단일 GPT 호출 래퍼"""
    res = client.chat.completions.create(
        model=MODEL_CHAT,
        messages=[{"role":"system","content":sys},{"role":"user","content":user}],
        response_format={"type":"json_object"} if format_json else None,
        temperature=0
    )
    return res.choices[0].message.content

async def agpt_many(prompt_tuples):
    """asyncio로 여러 GPT 호출 병렬"""
    loop = asyncio.get_event_loop()
    tasks=[]
    for sys,u in prompt_tuples:
        tasks.append(loop.run_in_executor(
            None, gpt_short_completion, sys, u, False))
    return await asyncio.gather(*tasks)

@st.cache_resource(show_spinner=False)
def embed_texts(texts:list[str])->np.ndarray:
    """OpenAI Embedding (cached)"""
    CHUNK=2048
    vectors=[]
    for i in range(0,len(texts),CHUNK):
        chunk=texts[i:i+CHUNK]
        vecs = client.embeddings.create(model=MODEL_EMB,input=chunk).data
        vectors.extend([v.embedding for v in vecs])
    return np.array(vectors)

# ──────────────────── CSV 로드 & 전처리 ──────────────────── #
def auto_read_csv(uploaded)->pd.DataFrame:
    for enc in ("utf-8","euc-kr","cp949"):
        try:
            return pd.read_csv(StringIO(uploaded.getvalue().decode(enc)))
        except Exception: continue
    raise UnicodeDecodeError("인코딩 자동 판독 실패")

# ────────────────── 텍스트 분석/토픽/감정 ────────────────── #
def basic_text_stats(series:pd.Series):
    texts=series.dropna().astype(str)
    if texts.empty: return None
    lens=texts.str.len()
    all_text=' '.join(texts)
    words=re.findall(r'[가-힣]+|[a-zA-Z]+|\d+',all_text.lower())
    words=[w for w in words if w not in STOPWORDS and len(w)>1]
    return {
        "total":len(texts),"avg":lens.mean(),"min":lens.min(),"max":lens.max(),
        "freq":Counter(words).most_common(20)
    }

def topic_clustering(texts:list[int], k:int=5):
    emb=embed_texts(texts)
    k=min(k,len(texts))
    model=KMeans(n_clusters=k,random_state=0,n_init="auto").fit(emb)
    return model.labels_, emb

# ──────────────────── PPT 자동 생성 ──────────────────── #
def make_ppt(overview:str, charts:list[tuple[str,bytes]])->str:
    prs=Presentation()
    slide=prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text="AI 설문 분석 요약"
    tf=slide.shapes.placeholders[1].text_frame
    for p in textwrap.wrap(overview,60):
        tf.add_paragraph().text=p
    # 차트 이미지 슬라이드
    for title,img_bytes in charts:
        s=prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text=title
        pic=s.shapes.add_picture(img_bytes,Inches(1),Inches(1.5),height=Inches(5))
    fname=f"survey_ai_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    prs.save(fname)
    return fname

# ──────────────────────── UI 시작 ──────────────────────── #
st.title("🧩 GPT 설문 데이터 슈트")

uploaded=st.file_uploader("CSV 업로드",type="csv")
if not uploaded: st.stop()

df=auto_read_csv(uploaded)
st.success(f"{len(df):,}행 · {len(df.columns)}열 로드 완료")
with st.expander("미리보기"):
    st.dataframe(df.head())

# -------- 1) 컬럼 타입 자동 추론 -------- #
sample=df.head(5).to_csv(index=False)
with st.spinner("GPT가 컬럼 의미 추론 중..."):
    gpt_types=gpt_infer_types(sample)

st.info("추론 결과를 확인/수정하세요")
col_types={}
two_cols=st.columns(2)
for i,c in enumerate(df.columns):
    with two_cols[i%2]:
        sel=st.selectbox(f"**{c}**",list(COLUMN_TYPES.keys()),
            index=list(COLUMN_TYPES.keys()).index(gpt_types.get(c,"other")),
            format_func=lambda x:COLUMN_TYPES[x],key=f"sel_{c}")
        col_types[c]=sel

st.divider()

# 확정 후 추가 기능 노출
if not st.button("🚀 타입 확정",type="primary"): st.stop()

st.session_state["df"]=df
st.session_state["types"]=col_types

# ====== 탭 구성 ====== #
tabs=st.tabs([
    "📊 개요 & 인사이트","🔬 텍스트 분석","🗣️ 자연어 Q&A",
    "👥 페르소나","🎯 설문 개선 피드백","📥 내보내기"
])

# ───── 2) 개요 & GPT 요약 ───── #
with tabs[0]:
    st.header("📊 전체 개요")
    total,len_q=len(df),len(df.columns)
    completion=(df.notna().sum().sum())/(total*len_q)*100
    st.metric("응답 수",f"{total:,}")
    st.metric("질문 수",len_q)
    st.metric("평균 응답률",f"{completion:.1f}%")
    # GPT 요약 버튼
    if st.toggle("AI 한 줄 요약 생성(토큰 사용)"):
        sys="당신은 통찰력 있는 데이터 분석가입니다. 아래 통계치를 한글로 2문장 요약."
        stats=f"응답수 {total}개, 질문 {len_q}개, 평균 응답률 {completion:.1f}%"
        with st.spinner("📜 요약 생성 중 ..."):
            summary=gpt_short_completion(sys,stats)
        st.success(summary)

# ───── 3) 텍스트 분석 ───── #
with tabs[1]:
    st.header("🔬 텍스트 질문 분석")
    text_cols=[c for c,t in col_types.items() if t.startswith("text_")]
    if not text_cols:
        st.info("텍스트형 질문이 없습니다"); st.stop()
    sel_col=st.selectbox("분석할 컬럼",text_cols)
    stats=basic_text_stats(df[sel_col])
    if not stats: st.warning("응답 없음"); st.stop()
    st.write(f"**응답 수** {stats['total']}  ·  **평균 길이** {stats['avg']:.0f}")
    # 토픽/감정 분석 토글
    run_topic=st.toggle("토픽/감정 분석 실행(토큰 사용)",key="tg_topic")
    if run_topic:
        texts=df[sel_col].dropna().astype(str).tolist()
        labels,emb=topic_clustering(texts,k= st.slider("클러스터 수",2,10,5))
        df_topic=pd.DataFrame({"text":texts,"topic":labels})
        st.dataframe(df_topic.head())
        col_counts=df_topic["topic"].value_counts().sort_index()
        fig=px.bar(x=col_counts.index,y=col_counts.values,labels={'x':'토픽','y':'개수'})
        st.plotly_chart(fig)
        # 토픽별 대표 문장 GPT 요약
        if st.toggle("토픽별 대표 키워드 요약"):
            prompts=[]
            for i in col_counts.index:
                sample=" | ".join(df_topic[df_topic.topic==i]["text"].head(5))
                sys="당신은 텍스트 분석가. 예시 문장을 보고 핵심 키워드 3개 한글로 콤마로."
                prompts.append((sys,sample[:4000]))
            with st.spinner("GPT 요약 중..."):
                summaries=asyncio.run(agpt_many(prompts))
            for i,sm in zip(col_counts.index,summaries):
                st.write(f"**토픽 {i}** → {sm}")

# ───── 4) 자연어 Q&A ───── #
with tabs[2]:
    st.header("🗣️ 자연어 질문 대시보드")
    q=st.chat_input("예: 남/여 비율은?")
    if "chat" not in st.session_state: st.session_state.chat=[]
    for role,m in st.session_state.chat:
        st.chat_message(role).write(m)
    if q:
        st.chat_message("user").write(q)
        st.session_state.chat.append(("user",q))
        # GPT에게 SQL 생성 지시
        cols=list(df.columns)
        sql=gpt_short_completion(
          "SQL expert. DataFrame columns: "+", ".join(cols)+
          ". Return pandasql SELECT for user's question. No explanation.",
          q)
        try:
            res=ps.sqldf(sql,locals())
            st.chat_message("assistant").write(res.head())
            st.session_state.chat.append(("assistant",res.head().to_markdown()))
        except Exception as e:
            st.chat_message("assistant").write(f"⚠️ 오류: {e}")
            st.session_state.chat.append(("assistant",str(e)))

# ───── 5) 응답자 페르소나 ───── #
with tabs[3]:
    st.header("👥 GPT 페르소나")
    pid_cols=[c for c,t in col_types.items() if t in ("email","phone","name")]
    if not pid_cols:
        st.info("개인 식별 컬럼이 없습니다"); st.stop()
    sel_pid=st.selectbox("기준 컬럼",pid_cols)
    emails=df[sel_pid].dropna().unique().tolist()[:50]
    if st.button("페르소나 생성 (최대 50명, 토큰 사용)"):
        prompts=[("너는 마케터. 다음 설문 응답으로 페르소나 한줄 요약.",
                  df[df[sel_pid]==e].iloc[0].to_json()) for e in emails]
        with st.spinner("GPT 생성 중"):
            pers=asyncio.run(agpt_many(prompts))
        st.dataframe(pd.DataFrame({"id":emails,"persona":pers}))

# ───── 6) 설문 개선 피드백 ───── #
with tabs[4]:
    st.header("🎯 설문 설계 개선 제안")
    if st.button("GPT 피드백 요청(토큰 사용)"):
        q_text="\n".join([f"{i+1}. {c}" for i,c in enumerate(df.columns)])
        sys="전문 리서처. 질문 리스트와 응답률 %.1f 를 보고 개선점 5가지."%completion
        fb=gpt_short_completion(sys,q_text)
        st.success(fb)

# ───── 7) 내보내기 (PPT 포함) ───── #
with tabs[5]:
    st.header("📥 데이터/보고서/PPT 내보내기")
    export=st.selectbox("형식 선택",["원본 CSV","익명 CSV","PPT 보고서"])
    if export=="원본 CSV":
        st.download_button("CSV 다운로드",
            df.to_csv(index=False).encode("utf-8-sig"),
            file_name="survey_raw.csv",mime="text/csv")
    elif export=="익명 CSV":
        anon=df.copy()
        for c,t in col_types.items():
            if t=="email": anon[c]=anon[c].str.replace(r'(.{2}).*@','\\1***@',regex=True)
            if t=="name":  anon[c]=anon[c].str.replace(r'(.).+','\\1*',regex=True)
        st.download_button("익명 CSV",
            anon.to_csv(index=False).encode("utf-8-sig"),
            file_name="survey_anon.csv",mime="text/csv")
    else:
        overview=f"응답 {total}개, 질문 {len_q}개, 응답률 {completion:.1f}%"
        st.info("차트 이미지를 넣으려면 상단 탭에서 '카메라 아이콘 → PNG'로 다운받은 후 업로드하세요.")
        imgs=st.file_uploader("PNG 차트 3개까지 (선택)",type=["png"],accept_multiple_files=True)
        charts=[]
        for f in imgs[:3]:
            charts.append((f.name,f.name))
            open(f.name,"wb").write(f.read())  # save temp
        fname=make_ppt(overview,charts)
        st.download_button("PPT 다운로드",open(fname,"rb").read(),
                           file_name=fname,mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
